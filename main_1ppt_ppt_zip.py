import re
import time
from docx import Document
from docx.shared import Pt
from docx.shared import Cm
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx2pdf import convert
from docx.shared import RGBColor  # 设置字体的颜色
from docx.oxml.ns import qn
from win32com import client as wc
from pptx import Presentation
import os
import shutil
import zipfile

zipfile.UNRAR_TOOL = "D:\\Program Files (x86)\\WinRAR\\UnRAR.exe"


def decompress_zip(zip_file_name, dir_name):
    """
    .zip 文件解压
    :param zip_file_name: zip 文件路径
    :param dir_name: 文件解压目录
    :return:
    """
    # 创建 zip 对象
    zip_obj = zipfile.ZipFile(zip_file_name)
    # 目录切换
    if not os.path.exists(dir_name):
        os.mkdir(dir_name)
    os.chdir(dir_name)
    # Extract all files into current directory.
    zip_obj.extractall()
    # zip_obj.extractall(dir_name)
    # 关闭
    zip_obj.close()


def get_slide_count(pptx_path):
    """
    获取PPTX文件的幻灯片数量（即页数）。

    参数:
    pptx_path (str): PPTX文件的路径。

    返回:
    int: 幻灯片的数量。
    """
    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)
    return slide_count


def replace_text_in_slide(slide, old_text, new_text):
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            if old_text in shape.text:
                shape.text = shape.text.replace(old_text, new_text)


if __name__ == '__main__':
    # 第一步：解压压缩包
    # zip_root_dir = "F:\\workspace\\www.1ppt.com\\2025-12-12\\www.zip_1ppt.com"
    # zip_dirs = sorted(os.listdir(zip_root_dir))
    # zip_files = sorted(os.listdir(zip_root_dir))
    # for zip_file in zip_files:
    #
    #     zip_file_path = zip_root_dir + "\\" + zip_file
    #     dst_file_path = "F:\\workspace\\www.1ppt.com\\2025-12-12\\www.uncompress_zip_1ppt.com"
    #     if not os.path.exists(dst_file_path):
    #         os.makedirs(dst_file_path)
    #     print("==========" + "开始解压" + zip_file_path + "==========")
    #     try:
    #         # 删除文件名中含有“图”字样文件
    #         if "图" in zip_file:
    #             os.remove(zip_file_path)
    #             continue
    #         # 删除文件名中含有“张”字样文件
    #         if "张" in zip_file:
    #             os.remove(zip_file_path)
    #             continue
    #         # 删除文件名中含有“套”字样文件
    #         if "套" in zip_file:
    #             os.remove(zip_file_path)
    #             continue
    #         # 删除文件名中含有“个”字样文件
    #         if "个" in zip_file:
    #             os.remove(zip_file_path)
    #             continue
    #         # 删除文件名中含有“页”字样文件
    #         if "页" in zip_file:
    #             os.remove(zip_file_path)
    #             continue
    #         # 删除文件名中含有“年”字样文件
    #         if "年" in zip_file:
    #             os.remove(zip_file_path)
    #             continue
    #         # 删除文件名中含有“素材”字样文件
    #         if "素材" in zip_file:
    #             os.remove(zip_file_path)
    #             continue
    #         decompress_zip(zip_file_path, dst_file_path+"\\"+zip_file.replace(".zip", "").replace("下载", "").replace("免费", ""))
    #     except Exception as e:
    #         print(e)
    #         continue
    #     print("==========" + "解压完成" + "==========")
    # exit()

    # 第二步：将文件夹中文件移出，并更改文件名称
    # root_dir = "F:\\workspace\\www.1ppt.com\\2025-12-12\\www.uncompress_zip_1ppt.com"
    # files = sorted(os.listdir(root_dir))
    # for file in files:
    #     file_path = root_dir + "\\" + file
    #     print(file_path)
    #
    #     # 查看一下是否是文件夹，如果是文件夹，则将文件移出
    #     if os.path.isdir(file_path):
    #         child_files = sorted(os.listdir(file_path))
    #         for child_file in child_files:
    #             extension = os.path.splitext(child_file)[-1]
    #             if extension not in [".pptx"]:
    #                 continue
    #             src_child_file_path = file_path + "\\" + child_file
    #             dst_child_file_path = root_dir + "\\" + file + extension
    #             try:
    #                 os.rename(src_child_file_path, dst_child_file_path)
    #             except WindowsError:
    #                 os.remove(dst_child_file_path)
    #                 os.rename(src_child_file_path, dst_child_file_path)
    # exit()

    # 第三步：替换幻灯片中文字
    root_dir = "F:\\workspace\\www.1ppt.com\\2025-12-12\\www.uncompress_zip_1ppt.com"
    finish_dir = "F:\\workspace\\www.1ppt.com\\2025-12-12\\www.finish_zip_1ppt.com"
    if not os.path.exists(finish_dir):
        os.makedirs(finish_dir)
    files = sorted(os.listdir(root_dir))
    for file in files:
        file_path = root_dir + "\\" + file
        print(file_path)

        finish_file_path = file_path.replace("uncompress", "finish")
        if not os.path.exists(finish_file_path):
            # 获取pptx文件总页数
            pptx_pages_count = get_slide_count(file_path)
            if pptx_pages_count <= 5:
                print(f"幻灯片数量（页数）: {pptx_pages_count}少于5，删除文件")
                os.remove(file_path)
                continue

            # 加载PPTX文件
            prs = Presentation(file_path)
            # 遍历所有幻灯片并替换文本
            for slide in prs.slides:
                # 替换文本
                replace_text_in_slide(slide, '第一PPT', 'XXXX')
                replace_text_in_slide(slide, '第一ppt', 'XXXX')
                replace_text_in_slide(slide, '1ppt', 'XXXX')
                replace_text_in_slide(slide, '1PPT', 'XXXX')
            # 删除最后一页
            rId = prs.slides._sldIdLst[-1].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[-1]
            # 保存修改后的PPTX文件
            prs.save(finish_file_path)
