import re
import time
from docx import Document
from docx.shared import Pt
from docx.shared import Cm
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx2pdf import convert
from docx.shared import RGBColor  # 设置字体的颜色
from docx.oxml.ns import qn
from pptx import Presentation
import os
import shutil


def get_slide_count(pptx_path):
    """
    获取PPTX文件的幻灯片数量（即页数）。

    参数:
    pptx_path (str): PPTX文件的路径。

    返回:
    int: 幻灯片的数量。
    """
    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)
    return slide_count


def replace_text_in_slide(slide, old_text, new_text):
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            if old_text in shape.text:
                shape.text = shape.text.replace(old_text, new_text)


if __name__ == '__main__':
    # 替换幻灯片中文字
    root_dir = "../www.officeplus.cn/2025-03-06/www.officeplus.cn/Ppt模板/"
    finish_dir = "../www.officeplus.cn/2025-03-06/ppt.officeplus.cn/"
    if not os.path.exists(finish_dir):
        os.makedirs(finish_dir)
    files = sorted(os.listdir(root_dir))
    for file in files:
        file_path = root_dir + file
        print(file_path)

        finish_file_path = finish_dir + file
        if not os.path.exists(finish_file_path):
            # 获取pptx文件总页数
            pptx_pages_count = get_slide_count(file_path)
            if pptx_pages_count <= 5:
                print(f"幻灯片数量（页数）: {pptx_pages_count}少于5，删除文件")
                os.remove(file_path)
                continue

            # 加载PPTX文件
            prs = Presentation(file_path)
            # 遍历所有幻灯片并替换文本
            for slide in prs.slides:
                # 替换文本
                replace_text_in_slide(slide, 'officeplus', 'XXXX')
                replace_text_in_slide(slide, 'OfficePlus', 'XXXX')
                replace_text_in_slide(slide, 'officePlus', 'XXXX')
                replace_text_in_slide(slide, 'OfficePLUS', 'XXXX')
            # 删除最后一页
            # rId = prs.slides._sldIdLst[-1].rId
            # prs.part.drop_rel(rId)
            # del prs.slides._sldIdLst[-1]
            # 保存修改后的PPTX文件
            prs.save(finish_file_path)
